"""
📄 Document Loader System
Day 25: Advanced document processing with multiple formats and sources
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Iterator
from dataclasses import dataclass, field
from datetime import datetime
import hashlib
from concurrent.futures import ThreadPoolExecutor, as_completed
from abc import ABC, abstractmethod

# Try to import required packages
try:
    import pypdf
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False
    print("⚠️ PyPDF not installed. Install with: pip install pypdf")

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️ python-docx not installed. Install with: pip install python-docx")

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️ python-pptx not installed. Install with: pip install python-pptx")

try:
    from unstructured.partition.auto import partition
    from unstructured.cleaners.core import clean_extra_whitespace
    HAS_UNSTRUCTURED = True
except ImportError:
    HAS_UNSTRUCTURED = False
    print("⚠️ Unstructured not installed. Install with: pip install unstructured")

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    print("⚠️ pandas not installed. Install with: pip install pandas")

try:
    import requests
    from bs4 import BeautifulSoup
    HAS_WEB = True
except ImportError:
    HAS_WEB = False
    print("⚠️ requests/beautifulsoup4 not installed. Install with: pip install requests beautifulsoup4")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

@dataclass
class Document:
    """Represents a loaded document with metadata"""
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    source: str = ""
    page_number: Optional[int] = None
    chunk_id: Optional[str] = None
    
    def __post_init__(self):
        """Initialize document with default metadata"""
        if 'id' not in self.metadata:
            self.metadata['id'] = hashlib.md5(self.content.encode()).hexdigest()[:16]
        if 'timestamp' not in self.metadata:
            self.metadata['timestamp'] = datetime.now().isoformat()
        if 'source' not in self.metadata and self.source:
            self.metadata['source'] = self.source
        if 'page_number' not in self.metadata and self.page_number is not None:
            self.metadata['page_number'] = self.page_number

class DocumentLoader(ABC):
    """Abstract base class for document loaders"""
    
    def __init__(self, file_path: Optional[str] = None, content: Optional[str] = None):
        self.file_path = file_path
        self.content = content
        self.documents: List[Document] = []
    
    @abstractmethod
    def load(self) -> List[Document]:
        """Load documents from source"""
        pass
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ""
        
        # Remove extra whitespace
        text = ' '.join(text.split())
        
        # Remove non-printable characters
        text = ''.join(char for char in text if char.isprintable() or char in '\n\r\t')
        
        return text.strip()
    
    def get_stats(self) -> Dict[str, Any]:
        """Get loading statistics"""
        total_chars = sum(len(doc.content) for doc in self.documents)
        total_words = sum(len(doc.content.split()) for doc in self.documents)
        
        return {
            'document_count': len(self.documents),
            'total_characters': total_chars,
            'total_words': total_words,
            'avg_words_per_doc': total_words / max(len(self.documents), 1),
            'file_path': self.file_path,
            'loader_type': self.__class__.__name__
        }

class PDFLoader(DocumentLoader):
    """PDF document loader using PyPDF"""
    
    def load(self) -> List[Document]:
        """Load PDF document"""
        if not HAS_PYPDF:
            raise ImportError("PyPDF not installed. Install with: pip install pypdf")
        
        if not self.file_path or not os.path.exists(self.file_path):
            raise FileNotFoundError(f"PDF file not found: {self.file_path}")
        
        logger.info(f"Loading PDF: {self.file_path}")
        
        try:
            reader = PdfReader(self.file_path)
            total_pages = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text:
                        cleaned_text = self.clean_text(text)
                        
                        document = Document(
                            content=cleaned_text,
                            metadata={
                                'source': self.file_path,
                                'page_number': page_num,
                                'total_pages': total_pages,
                                'format': 'pdf',
                                'file_size': os.path.getsize(self.file_path),
                                'title': reader.metadata.title if reader.metadata else None,
                                'author': reader.metadata.author if reader.metadata else None
                            },
                            source=self.file_path,
                            page_number=page_num
                        )
                        self.documents.append(document)
                        
                        logger.debug(f"Loaded page {page_num}/{total_pages}: {len(cleaned_text)} chars")
                
                except Exception as e:
                    logger.error(f"Error extracting text from page {page_num}: {e}")
            
            logger.info(f"Loaded {len(self.documents)} pages from PDF")
            return self.documents
        
        except Exception as e:
            logger.error(f"Error loading PDF {self.file_path}: {e}")
            raise

class DOCXLoader(DocumentLoader):
    """Microsoft Word document loader"""
    
    def load(self) -> List[Document]:
        """Load DOCX document"""
        if not HAS_DOCX:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        if not self.file_path or not os.path.exists(self.file_path):
            raise FileNotFoundError(f"DOCX file not found: {self.file_path}")
        
        logger.info(f"Loading DOCX: {self.file_path}")
        
        try:
            doc = docx.Document(self.file_path)
            
            # Extract text from paragraphs
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            full_text.append(cell.text)
            
            content = '\n\n'.join(full_text)
            cleaned_content = self.clean_text(content)
            
            if cleaned_content:
                document = Document(
                    content=cleaned_content,
                    metadata={
                        'source': self.file_path,
                        'format': 'docx',
                        'file_size': os.path.getsize(self.file_path),
                        'paragraph_count': len(doc.paragraphs),
                        'table_count': len(doc.tables)
                    },
                    source=self.file_path
                )
                self.documents.append(document)
            
            logger.info(f"Loaded DOCX: {len(cleaned_content)} characters")
            return self.documents
        
        except Exception as e:
            logger.error(f"Error loading DOCX {self.file_path}: {e}")
            raise

class PPTXLoader(DocumentLoader):
    """PowerPoint presentation loader"""
    
    def load(self) -> List[Document]:
        """Load PPTX presentation"""
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        if not self.file_path or not os.path.exists(self.file_path):
            raise FileNotFoundError(f"PPTX file not found: {self.file_path}")
        
        logger.info(f"Loading PPTX: {self.file_path}")
        
        try:
            presentation = pptx.Presentation(self.file_path)
            
            full_text = []
            slide_number = 0
            
            for slide in presentation.slides:
                slide_number += 1
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slide_content = f"Slide {slide_number}:\n" + '\n'.join(slide_text)
                    full_text.append(slide_content)
            
            content = '\n\n'.join(full_text)
            cleaned_content = self.clean_text(content)
            
            if cleaned_content:
                document = Document(
                    content=cleaned_content,
                    metadata={
                        'source': self.file_path,
                        'format': 'pptx',
                        'file_size': os.path.getsize(self.file_path),
                        'slide_count': slide_number,
                        'total_slides': len(presentation.slides)
                    },
                    source=self.file_path
                )
                self.documents.append(document)
            
            logger.info(f"Loaded PPTX: {len(presentation.slides)} slides, {len(cleaned_content)} characters")
            return self.documents
        
        except Exception as e:
            logger.error(f"Error loading PPTX {self.file_path}: {e}")
            raise

class TXTLoader(DocumentLoader):
    """Plain text file loader"""
    
    def load(self) -> List[Document]:
        """Load text file"""
        if not self.file_path or not os.path.exists(self.file_path):
            raise FileNotFoundError(f"Text file not found: {self.file_path}")
        
        logger.info(f"Loading text file: {self.file_path}")
        
        try:
            with open(self.file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            cleaned_content = self.clean_text(content)
            
            if cleaned_content:
                document = Document(
                    content=cleaned_content,
                    metadata={
                        'source': self.file_path,
                        'format': 'txt',
                        'file_size': os.path.getsize(self.file_path),
                        'encoding': 'utf-8'
                    },
                    source=self.file_path
                )
                self.documents.append(document)
            
            logger.info(f"Loaded text file: {len(cleaned_content)} characters")
            return self.documents
        
        except UnicodeDecodeError:
            # Try different encodings
            encodings = ['latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(self.file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    
                    cleaned_content = self.clean_text(content)
                    
                    if cleaned_content:
                        document = Document(
                            content=cleaned_content,
                            metadata={
                                'source': self.file_path,
                                'format': 'txt',
                                'file_size': os.path.getsize(self.file_path),
                                'encoding': encoding
                            },
                            source=self.file_path
                        )
                        self.documents.append(document)
                    
                    logger.info(f"Loaded text file with encoding {encoding}: {len(cleaned_content)} characters")
                    return self.documents
                
                except UnicodeDecodeError:
                    continue
            
            raise ValueError(f"Could not decode file {self.file_path} with any supported encoding")
        
        except Exception as e:
            logger.error(f"Error loading text file {self.file_path}: {e}")
            raise

class CSVLoader(DocumentLoader):
    """CSV file loader with pandas"""
    
    def __init__(self, file_path: Optional[str] = None, content: Optional[str] = None,
                 delimiter: str = ',', max_rows: int = 1000):
        super().__init__(file_path, content)
        self.delimiter = delimiter
        self.max_rows = max_rows
    
    def load(self) -> List[Document]:
        """Load CSV file"""
        if not HAS_PANDAS:
            raise ImportError("pandas not installed. Install with: pip install pandas")
        
        if not self.file_path or not os.path.exists(self.file_path):
            raise FileNotFoundError(f"CSV file not found: {self.file_path}")
        
        logger.info(f"Loading CSV: {self.file_path}")
        
        try:
            # Read CSV
            df = pd.read_csv(self.file_path, delimiter=self.delimiter, nrows=self.max_rows)
            
            # Convert to text representation
            content_parts = []
            
            # Add column names
            columns = ', '.join(df.columns.tolist())
            content_parts.append(f"Columns: {columns}\n")
            
            # Add sample rows
            for idx, row in df.head(10).iterrows():
                row_text = ', '.join([str(val) for val in row.tolist()])
                content_parts.append(f"Row {idx}: {row_text}")
            
            content = '\n'.join(content_parts)
            cleaned_content = self.clean_text(content)
            
            if cleaned_content:
                document = Document(
                    content=cleaned_content,
                    metadata={
                        'source': self.file_path,
                        'format': 'csv',
                        'file_size': os.path.getsize(self.file_path),
                        'rows': len(df),
                        'columns': len(df.columns),
                        'delimiter': self.delimiter
                    },
                    source=self.file_path
                )
                self.documents.append(document)
            
            logger.info(f"Loaded CSV: {len(df)} rows, {len(df.columns)} columns")
            return self.documents
        
        except Exception as e:
            logger.error(f"Error loading CSV {self.file_path}: {e}")
            raise

class WebLoader(DocumentLoader):
    """Web page loader"""
    
    def __init__(self, url: str, timeout: int = 10):
        super().__init__()
        self.url = url
        self.timeout = timeout
    
    def load(self) -> List[Document]:
        """Load web page content"""
        if not HAS_WEB:
            raise ImportError("requests/beautifulsoup4 not installed")
        
        logger.info(f"Loading web page: {self.url}")
        
        try:
            # Fetch web page
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(self.url, headers=headers, timeout=self.timeout)
            response.raise_for_status()
            
            # Parse HTML
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            
            # Extract text
            text = soup.get_text()
            
            # Clean text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            cleaned_content = self.clean_text(text)
            
            if cleaned_content:
                document = Document(
                    content=cleaned_content,
                    metadata={
                        'source': self.url,
                        'format': 'web',
                        'status_code': response.status_code,
                        'content_type': response.headers.get('content-type'),
                        'title': soup.title.string if soup.title else None,
                        'timestamp': datetime.now().isoformat()
                    },
                    source=self.url
                )
                self.documents.append(document)
            
            logger.info(f"Loaded web page: {len(cleaned_content)} characters")
            return self.documents
        
        except Exception as e:
            logger.error(f"Error loading web page {self.url}: {e}")
            raise

class UnstructuredLoader(DocumentLoader):
    """Advanced document loader using Unstructured library"""
    
    def __init__(self, file_path: Optional[str] = None, content: Optional[str] = None,
                 strategy: str = "auto", languages: List[str] = None):
        super().__init__(file_path, content)
        self.strategy = strategy
        self.languages = languages or ["eng"]
    
    def load(self) -> List[Document]:
        """Load document using Unstructured"""
        if not HAS_UNSTRUCTURED:
            raise ImportError("Unstructured not installed")
        
        if not self.file_path or not os.path.exists(self.file_path):
            raise FileNotFoundError(f"File not found: {self.file_path}")
        
        logger.info(f"Loading with Unstructured: {self.file_path}")
        
        try:
            # Partition document
            elements = partition(
                filename=self.file_path,
                strategy=self.strategy,
                languages=self.languages
            )
            
            # Group elements by page if available
            pages = {}
            for element in elements:
                page_number = element.metadata.page_number if hasattr(element, 'metadata') else 0
                if page_number not in pages:
                    pages[page_number] = []
                pages[page_number].append(str(element))
            
            # Create documents for each page
            for page_num, page_elements in pages.items():
                content = '\n'.join(page_elements)
                cleaned_content = self.clean_text(content)
                
                if cleaned_content:
                    document = Document(
                        content=cleaned_content,
                        metadata={
                            'source': self.file_path,
                            'format': Path(self.file_path).suffix[1:],
                            'file_size': os.path.getsize(self.file_path),
                            'page_number': page_num,
                            'loader': 'unstructured',
                            'strategy': self.strategy,
                            'languages': self.languages
                        },
                        source=self.file_path,
                        page_number=page_num
                    )
                    self.documents.append(document)
            
            logger.info(f"Loaded {len(self.documents)} pages with Unstructured")
            return self.documents
        
        except Exception as e:
            logger.error(f"Error loading with Unstructured {self.file_path}: {e}")
            raise

class DirectoryLoader:
    """Load documents from a directory"""
    
    def __init__(self, directory_path: str, recursive: bool = True,
                 file_extensions: Optional[List[str]] = None, max_workers: int = 4):
        self.directory_path = Path(directory_path)
        self.recursive = recursive
        self.file_extensions = file_extensions or ['.pdf', '.txt', '.docx', '.pptx', '.csv']
        self.max_workers = max_workers
        self.loaded_documents: List[Document] = []
        self.failed_files: List[Dict[str, Any]] = []
    
    def get_file_loader(self, file_path: Path) -> Optional[DocumentLoader]:
        """Get appropriate loader for file extension"""
        extension = file_path.suffix.lower()
        
        loader_map = {
            '.pdf': PDFLoader,
            '.txt': TXTLoader,
            '.docx': DOCXLoader,
            '.pptx': PPTXLoader,
            '.csv': CSVLoader,
        }
        
        if extension in loader_map:
            return loader_map[extension](str(file_path))
        
        # Try Unstructured for other file types
        if HAS_UNSTRUCTURED:
            return UnstructuredLoader(str(file_path))
        
        return None
    
    def load(self) -> List[Document]:
        """Load all documents from directory"""
        if not self.directory_path.exists():
            raise FileNotFoundError(f"Directory not found: {self.directory_path}")
        
        logger.info(f"Loading documents from: {self.directory_path}")
        
        # Collect all files
        files = []
        pattern = "**/*" if self.recursive else "*"
        
        for file_path in self.directory_path.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in self.file_extensions:
                files.append(file_path)
        
        logger.info(f"Found {len(files)} files to process")
        
        # Process files in parallel
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            future_to_file = {}
            
            for file_path in files:
                loader = self.get_file_loader(file_path)
                if loader:
                    future = executor.submit(self._load_file, loader, file_path)
                    future_to_file[future] = file_path
            
            # Collect results
            for future in as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    documents = future.result(timeout=60)
                    self.loaded_documents.extend(documents)
                    logger.info(f"✓ Loaded {len(documents)} documents from {file_path.name}")
                except Exception as e:
                    self.failed_files.append({
                        'file': str(file_path),
                        'error': str(e)
                    })
                    logger.error(f"✗ Failed to load {file_path.name}: {e}")
        
        logger.info(f"Loaded {len(self.loaded_documents)} total documents")
        logger.info(f"Failed to load {len(self.failed_files)} files")
        
        return self.loaded_documents
    
    def _load_file(self, loader: DocumentLoader, file_path: Path) -> List[Document]:
        """Load a single file"""
        try:
            return loader.load()
        except Exception as e:
            raise Exception(f"Error loading {file_path}: {e}")
    
    def get_stats(self) -> Dict[str, Any]:
        """Get directory loading statistics"""
        total_chars = sum(len(doc.content) for doc in self.loaded_documents)
        total_words = sum(len(doc.content.split()) for doc in self.loaded_documents)
        
        # Count by format
        formats = {}
        for doc in self.loaded_documents:
            fmt = doc.metadata.get('format', 'unknown')
            formats[fmt] = formats.get(fmt, 0) + 1
        
        return {
            'total_documents': len(self.loaded_documents),
            'total_characters': total_chars,
            'total_words': total_words,
            'failed_files': len(self.failed_files),
            'formats': formats,
            'directory': str(self.directory_path)
        }

class DocumentLoaderFactory:
    """Factory for creating document loaders"""
    
    @staticmethod
    def create_loader(file_path: str) -> DocumentLoader:
        """Create appropriate loader based on file extension"""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        loader_map = {
            '.pdf': PDFLoader,
            '.txt': TXTLoader,
            '.docx': DOCXLoader,
            '.pptx': PPTXLoader,
            '.csv': CSVLoader,
        }
        
        if extension in loader_map:
            return loader_map[extension](file_path)
        
        # Try Unstructured for other formats
        if HAS_UNSTRUCTURED:
            return UnstructuredLoader(file_path)
        
        raise ValueError(f"Unsupported file format: {extension}")
    
    @staticmethod
    def create_web_loader(url: str) -> WebLoader:
        """Create web page loader"""
        return WebLoader(url)
    
    @staticmethod
    def create_directory_loader(directory_path: str, **kwargs) -> DirectoryLoader:
        """Create directory loader"""
        return DirectoryLoader(directory_path, **kwargs)

def demo_document_loading():
    """Demonstrate document loading capabilities"""
    print("=" * 70)
    print("📄 DOCUMENT LOADER DEMONSTRATION (Day 25)")
    print("=" * 70)
    
    # Create a temporary directory for demo
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        print(f"\n📁 Created temporary directory: {temp_path}")
        
        # Create sample files
        print("\n1️⃣ Creating sample documents...")
        
        # Sample text file
        txt_file = temp_path / "sample.txt"
        txt_content = """Artificial Intelligence (AI) is transforming industries across the globe.
Machine learning algorithms enable computers to learn from data and make predictions.
Natural Language Processing (NLP) allows computers to understand human language.
Computer vision enables machines to interpret visual information."""
        txt_file.write_text(txt_content)
        print(f"   Created: {txt_file.name}")
        
        # Sample CSV file
        csv_file = temp_path / "sample.csv"
        csv_content = """name,age,city,occupation
Alice,30,New York,Data Scientist
Bob,25,San Francisco,Software Engineer
Charlie,35,London,AI Researcher
Diana,28,Boston,Machine Learning Engineer"""
        csv_file.write_text(csv_content)
        print(f"   Created: {csv_file.name}")
        
        # Test individual loaders
        print("\n2️⃣ Testing individual loaders...")
        
        # Test TXT loader
        try:
            txt_loader = TXTLoader(str(txt_file))
            txt_docs = txt_loader.load()
            print(f"   ✓ TXT Loader: {len(txt_docs)} document, {len(txt_docs[0].content) if txt_docs else 0} chars")
        except Exception as e:
            print(f"   ✗ TXT Loader failed: {e}")
        
        # Test CSV loader
        try:
            if HAS_PANDAS:
                csv_loader = CSVLoader(str(csv_file))
                csv_docs = csv_loader.load()
                print(f"   ✓ CSV Loader: {len(csv_docs)} document")
            else:
                print("   ⚠️ CSV Loader skipped (pandas not installed)")
        except Exception as e:
            print(f"   ✗ CSV Loader failed: {e}")
        
        # Test factory pattern
        print("\n3️⃣ Testing loader factory...")
        
        try:
            factory = DocumentLoaderFactory()
            
            # Get TXT loader from factory
            txt_loader = factory.create_loader(str(txt_file))
            print(f"   ✓ Created TXT loader via factory: {txt_loader.__class__.__name__}")
            
            # Get CSV loader from factory
            csv_loader = factory.create_loader(str(csv_file))
            print(f"   ✓ Created CSV loader via factory: {csv_loader.__class__.__name__}")
        
        except Exception as e:
            print(f"   ✗ Factory failed: {e}")
        
        # Test directory loader
        print("\n4️⃣ Testing directory loader...")
        
        try:
            dir_loader = DirectoryLoader(
                str(temp_path),
                recursive=False,
                max_workers=2
            )
            
            dir_docs = dir_loader.load()
            stats = dir_loader.get_stats()
            
            print(f"   ✓ Directory Loader: {stats['total_documents']} documents")
            print(f"     Formats: {stats['formats']}")
            print(f"     Total words: {stats['total_words']}")
            
            # Show document metadata
            if dir_docs:
                print(f"\n   📋 Sample document metadata:")
                for i, doc in enumerate(dir_docs[:2], 1):
                    print(f"     {i}. Source: {doc.metadata.get('source', 'N/A')}")
                    print(f"        Format: {doc.metadata.get('format', 'N/A')}")
                    print(f"        Words: {len(doc.content.split())}")
        
        except Exception as e:
            print(f"   ✗ Directory Loader failed: {e}")
        
        # Test web loader (if available)
        if HAS_WEB:
            print("\n5️⃣ Testing web loader (example URL)...")
            
            try:
                # Use a simple example URL
                test_url = "http://example.com"
                web_loader = WebLoader(test_url)
                web_docs = web_loader.load()
                
                if web_docs:
                    print(f"   ✓ Web Loader: {len(web_docs[0].content)} characters loaded")
                    print(f"     Title: {web_docs[0].metadata.get('title', 'N/A')}")
                else:
                    print("   ⚠️ Web Loader: No content loaded")
            
            except Exception as e:
                print(f"   ✗ Web Loader failed: {e}")
        else:
            print("\n5️⃣ Web loader test skipped (requests/beautifulsoup4 not installed)")
    
    print("\n" + "=" * 70)
    print("✅ Document Loader Demonstration Complete!")
    print("📚 Next: Text Splitting Strategies (Day 26)")

def benchmark_loaders():
    """Benchmark different document loaders"""
    print("\n⏱️ Document Loader Benchmark")
    print("=" * 60)
    
    import time
    from collections import defaultdict
    
    # Create test files
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        # Create files of different sizes
        file_sizes = [1000, 10000, 50000]  # Characters
        results = defaultdict(list)
        
        for size in file_sizes:
            print(f"\n📊 Testing with {size:,} character files:")
            
            # Create text file
            txt_file = temp_path / f"test_{size}.txt"
            content = "X" * size
            txt_file.write_text(content)
            
            # Test TXT loader
            try:
                start = time.time()
                loader = TXTLoader(str(txt_file))
                docs = loader.load()
                elapsed = time.time() - start
                
                results['txt'].append((size, elapsed))
                print(f"   TXT Loader: {elapsed:.3f}s")
            except Exception as e:
                print(f"   TXT Loader failed: {e}")
        
        # Print summary
        print("\n📈 Benchmark Summary:")
        print("-" * 40)
        print(f"{'Format':<10} {'Size':<12} {'Time (s)':<10} {'Chars/s':<12}")
        print("-" * 40)
        
        for format_name, format_results in results.items():
            for size, elapsed in format_results:
                chars_per_sec = size / elapsed if elapsed > 0 else 0
                print(f"{format_name:<10} {size:<12,} {elapsed:<10.3f} {chars_per_sec:<12,.0f}")

def interactive_loader_demo():
    """Interactive document loader demonstration"""
    print("\n🎮 Interactive Document Loader")
    print("=" * 60)
    
    print("\n📂 Available commands:")
    print("  /load [file_path] - Load a document")
    print("  /dir [directory] - Load all documents from directory")
    print("  /web [url] - Load web page")
    print("  /stats - Show loading statistics")
    print("  /preview [n] - Preview first n characters of loaded documents")
    print("  /clear - Clear loaded documents")
    print("  /quit - Exit")
    
    loaded_documents = []
    
    while True:
        try:
            command = input("\n> ").strip()
            
            if command.lower() == '/quit':
                print("Goodbye!")
                break
            
            elif command.lower() == '/stats':
                if not loaded_documents:
                    print("❌ No documents loaded")
                else:
                    total_chars = sum(len(doc.content) for doc in loaded_documents)
                    total_words = sum(len(doc.content.split()) for doc in loaded_documents)
                    
                    print(f"\n📊 Loading Statistics:")
                    print(f"  Total documents: {len(loaded_documents)}")
                    print(f"  Total characters: {total_chars:,}")
                    print(f"  Total words: {total_words:,}")
                    print(f"  Average words per doc: {total_words/len(loaded_documents):.0f}")
                    
                    # Count by format
                    formats = {}
                    for doc in loaded_documents:
                        fmt = doc.metadata.get('format', 'unknown')
                        formats[fmt] = formats.get(fmt, 0) + 1
                    
                    print(f"\n  Documents by format:")
                    for fmt, count in formats.items():
                        print(f"    {fmt}: {count}")
            
            elif command.lower() == '/clear':
                loaded_documents.clear()
                print("✅ Cleared all loaded documents")
            
            elif command.lower().startswith('/preview '):
                try:
                    n = int(command.split()[1])
                    if not loaded_documents:
                        print("❌ No documents loaded")
                    else:
                        print(f"\n📄 Preview (first {n} characters):")
                        for i, doc in enumerate(loaded_documents[:3], 1):
                            print(f"\n{i}. Source: {doc.metadata.get('source', 'N/A')}")
                            print(f"   Content: {doc.content[:n]}...")
                except (IndexError, ValueError):
                    print("❌ Usage: /preview [number_of_characters]")
            
            elif command.lower().startswith('/load '):
                try:
                    file_path = command[6:].strip()
                    if not os.path.exists(file_path):
                        print(f"❌ File not found: {file_path}")
                        continue
                    
                    print(f"📥 Loading: {file_path}")
                    
                    # Use factory to create loader
                    loader = DocumentLoaderFactory.create_loader(file_path)
                    documents = loader.load()
                    
                    loaded_documents.extend(documents)
                    print(f"✅ Loaded {len(documents)} document(s)")
                    
                    # Show document info
                    for doc in documents:
                        print(f"   - {len(doc.content):,} characters, format: {doc.metadata.get('format', 'unknown')}")
                
                except Exception as e:
                    print(f"❌ Error loading document: {e}")
            
            elif command.lower().startswith('/dir '):
                try:
                    dir_path = command[5:].strip()
                    if not os.path.exists(dir_path):
                        print(f"❌ Directory not found: {dir_path}")
                        continue
                    
                    print(f"📂 Loading directory: {dir_path}")
                    
                    loader = DirectoryLoader(dir_path, max_workers=4)
                    documents = loader.load()
                    
                    loaded_documents.extend(documents)
                    print(f"✅ Loaded {len(documents)} documents from directory")
                    
                    # Show stats
                    stats = loader.get_stats()
                    print(f"   Formats: {stats['formats']}")
                    print(f"   Failed files: {stats['failed_files']}")
                
                except Exception as e:
                    print(f"❌ Error loading directory: {e}")
            
            elif command.lower().startswith('/web '):
                try:
                    if not HAS_WEB:
                        print("❌ Web loading requires requests and beautifulsoup4")
                        continue
                    
                    url = command[5:].strip()
                    print(f"🌐 Loading web page: {url}")
                    
                    loader = WebLoader(url)
                    documents = loader.load()
                    
                    loaded_documents.extend(documents)
                    print(f"✅ Loaded web page: {len(documents[0].content):,} characters")
                    print(f"   Title: {documents[0].metadata.get('title', 'N/A')}")
                
                except Exception as e:
                    print(f"❌ Error loading web page: {e}")
            
            else:
                print("❌ Unknown command. Type /help for commands.")
        
        except KeyboardInterrupt:
            print("\n\nExiting...")
            break
        except Exception as e:
            print(f"❌ Error: {e}")

def main():
    """Run document loader demonstrations"""
    print("🚀 Advanced Document Loading System")
    print("Day 25: Multi-format Document Processing with Parallel Loading")
    
    # Demo 1: Document loading demonstration
    demo_document_loading()
    
    # Demo 2: Benchmarking
    benchmark_loaders()
    
    # Demo 3: Interactive demo
    try:
        run_interactive = input("\n🎮 Run interactive loader demo? (y/n): ").lower()
        if run_interactive == 'y':
            interactive_loader_demo()
    except:
        print("Skipping interactive demo...")
    
    print("\n" + "=" * 70)
    print("✅ Day 25 Complete!")
    print("📚 You've built a production-ready document loading system!")
    print("   Next: Text Splitting Strategies (Day 26)")

if __name__ == "__main__":
    main()